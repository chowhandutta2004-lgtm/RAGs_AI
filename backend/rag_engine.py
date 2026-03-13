import os
import io
import re
import json
import uuid
import csv
from collections import Counter
from datetime import datetime, timezone
from typing import List, Dict, Any, Optional, AsyncGenerator

from dotenv import load_dotenv

# Document loaders
import pypdf
import docx

# LangChain
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_chroma import Chroma
from langchain_core.messages import HumanMessage, AIMessage

load_dotenv()

STOP_WORDS = {
    'the', 'and', 'for', 'are', 'but', 'not', 'you', 'all', 'can', 'her', 'was', 'one',
    'our', 'out', 'day', 'get', 'has', 'him', 'his', 'how', 'its', 'may', 'new', 'now',
    'old', 'see', 'two', 'way', 'who', 'boy', 'did', 'man', 'men', 'own', 'say', 'she',
    'too', 'use', 'with', 'from', 'this', 'that', 'they', 'have', 'been', 'were', 'said',
    'each', 'which', 'their', 'will', 'when', 'more', 'what', 'some', 'than', 'into',
    'your', 'also', 'just', 'time', 'like', 'then', 'them', 'over', 'such', 'only',
    'come', 'could', 'there', 'would', 'other', 'about', 'these', 'many', 'well', 'much',
    'any', 'after', 'most', 'very', 'even', 'back', 'good', 'know', 'take', 'make',
    'does', 'part', 'need', 'same', 'long', 'down', 'both', 'here', 'through', 'think',
    'first', 'those', 'being', 'where', 'while', 'should', 'every', 'found', 'still',
    'between', 'without', 'during', 'before', 'under', 'never', 'always', 'made', 'give',
    'great', 'since', 'right', 'things', 'place', 'world', 'again', 'might', 'point',
    'large', 'general', 'however', 'often', 'show', 'going', 'used', 'small', 'number',
    'another', 'until', 'less', 'across', 'once', 'given', 'among', 'within', 'including',
    'along', 'himself', 'themselves', 'because', 'against', 'important', 'something',
    'whether', 'around', 'several', 'example', 'able', 'became', 'above', 'though',
    'either', 'together', 'already', 'later', 'although', 'years', 'people', 'work',
    'system', 'order', 'process', 'level', 'type', 'provide', 'following', 'total',
    'further', 'based', 'available', 'must', 'called', 'possible', 'include', 'result',
    'case', 'form', 'local', 'name', 'high', 'hand', 'seem', 'turn',
}


class RAGEngine:
    def __init__(self):
        self.openai_api_key = os.getenv("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError("OPENAI_API_KEY not found in .env file!")

        self.embeddings = OpenAIEmbeddings(
            openai_api_key=self.openai_api_key,
            model="text-embedding-3-small"
        )

        self.llm = ChatOpenAI(
            openai_api_key=self.openai_api_key,
            model_name="gpt-4o",
            temperature=0.2
        )

        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " "]
        )

        self.persist_dir = "./chroma_db"
        self.user_data_dir = "./user_data"

    def _get_vectorstore(self, user_id: str) -> Chroma:
        return Chroma(
            collection_name=f"user_{user_id}",
            persist_directory=self.persist_dir,
            embedding_function=self.embeddings
        )

    def _load_metadata(self, user_id: str) -> list:
        path = os.path.join(self.user_data_dir, user_id, "metadata.json")
        if os.path.exists(path):
            with open(path, "r") as f:
                return json.load(f)
        return []

    def _save_metadata(self, user_id: str, data: list):
        dir_path = os.path.join(self.user_data_dir, user_id)
        os.makedirs(dir_path, exist_ok=True)
        path = os.path.join(dir_path, "metadata.json")
        with open(path, "w") as f:
            json.dump(data, f)

    # ── Session management ────────────────────────────────────────────────────

    def _load_sessions(self, user_id: str) -> list:
        path = os.path.join(self.user_data_dir, user_id, "sessions.json")
        if os.path.exists(path):
            with open(path, "r") as f:
                return json.load(f)
        return []

    def _save_sessions(self, user_id: str, data: list):
        dir_path = os.path.join(self.user_data_dir, user_id)
        os.makedirs(dir_path, exist_ok=True)
        path = os.path.join(dir_path, "sessions.json")
        with open(path, "w") as f:
            json.dump(data, f)

    def get_sessions(self, user_id: str) -> list:
        sessions = self._load_sessions(user_id)
        # Return summary (no messages payload) sorted newest-first
        return [
            {
                "id": s["id"],
                "name": s["name"],
                "created_at": s["created_at"],
                "message_count": len(s.get("messages", []))
            }
            for s in reversed(sessions)
        ]

    def create_session(self, user_id: str) -> dict:
        sessions = self._load_sessions(user_id)
        session = {
            "id": str(uuid.uuid4()),
            "name": "New Chat",
            "created_at": datetime.now(timezone.utc).isoformat(),
            "messages": []
        }
        sessions.append(session)
        self._save_sessions(user_id, sessions)
        return {"id": session["id"], "name": session["name"], "created_at": session["created_at"]}

    def get_session_messages(self, session_id: str, user_id: str) -> list:
        sessions = self._load_sessions(user_id)
        for s in sessions:
            if s["id"] == session_id:
                return s.get("messages", [])
        return []

    def save_session_exchange(self, session_id: str, user_id: str, question: str, answer: str, sources: list, confidence, used_context: bool):
        sessions = self._load_sessions(user_id)
        for s in sessions:
            if s["id"] == session_id:
                if not s.get("messages"):
                    # Auto-name from first user message
                    s["name"] = question[:40].strip()
                s.setdefault("messages", [])
                s["messages"].append({"role": "user", "content": question})
                s["messages"].append({
                    "role": "assistant",
                    "content": answer,
                    "sources": sources,
                    "confidence": confidence,
                    "used_context": used_context
                })
                break
        self._save_sessions(user_id, sessions)

    def delete_session(self, session_id: str, user_id: str):
        sessions = self._load_sessions(user_id)
        sessions = [s for s in sessions if s["id"] != session_id]
        self._save_sessions(user_id, sessions)

    # ── Keyword extraction ────────────────────────────────────────────────────

    def _extract_keywords(self, text: str, top_n: int = 10) -> list:
        words = re.findall(r'\b[a-z]{4,}\b', text.lower())
        filtered = [w for w in words if w not in STOP_WORDS]
        counter = Counter(filtered)
        return [{"word": word, "count": count} for word, count in counter.most_common(top_n)]

    # ── Text extraction ───────────────────────────────────────────────────────

    def extract_text(self, contents: bytes, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()

        if ext == '.pdf':
            reader = pypdf.PdfReader(io.BytesIO(contents))
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text

        elif ext == '.docx':
            doc = docx.Document(io.BytesIO(contents))
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext in ('.txt', '.md'):
            return contents.decode('utf-8', errors='ignore')

        elif ext == '.csv':
            text = ""
            reader = csv.reader(io.StringIO(contents.decode('utf-8', errors='ignore')))
            for row in reader:
                text += ", ".join(row) + "\n"
            return text

        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(contents), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        lines.append(row_text)
            return "\n".join(lines)

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(io.BytesIO(contents))
            lines = []
            for slide_num, slide in enumerate(prs.slides, 1):
                lines.append(f"[Slide {slide_num}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)

        raise ValueError(f"Unsupported file type: {ext}")

    # ── Ingestion ─────────────────────────────────────────────────────────────

    def ingest_document(self, contents: bytes, filename: str, user_id: str) -> Dict:
        existing = self._load_metadata(user_id)
        if any(d["name"] == filename for d in existing):
            self.delete_document(filename, user_id)
        text = self.extract_text(contents, filename)
        word_count = len(text.split())
        keywords = self._extract_keywords(text)

        chunks = self.splitter.create_documents(
            texts=[text],
            metadatas=[{"source": filename, "filename": filename}]
        )

        vectorstore = self._get_vectorstore(user_id)
        vectorstore.add_documents(chunks)

        metadata = self._load_metadata(user_id)
        metadata.append({
            "name": filename,
            "chunks": len(chunks),
            "words": word_count,
            "type": os.path.splitext(filename)[1].upper().replace(".", ""),
            "keywords": keywords,
            "uploaded_at": datetime.now(timezone.utc).isoformat()
        })
        self._save_metadata(user_id, metadata)

        return {"chunks": len(chunks), "words": word_count}

    def ingest_url(self, url: str, user_id: str) -> Dict:
        import httpx
        from bs4 import BeautifulSoup
        from urllib.parse import urlparse

        with httpx.Client(follow_redirects=True, timeout=20) as client:
            response = client.get(url, headers={"User-Agent": "RAGs_AI/1.0"})
            response.raise_for_status()
            html = response.text

        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)

        parsed = urlparse(url)
        raw_name = (parsed.netloc + parsed.path).strip("/").replace("/", "_")
        # Sanitize: keep alphanumerics, dots, dashes, underscores
        safe_name = re.sub(r'[^\w.\-]', '_', raw_name)[:80] + ".url"

        word_count = len(text.split())
        keywords = self._extract_keywords(text)

        chunks = self.splitter.create_documents(
            texts=[text],
            metadatas=[{"source": safe_name, "filename": safe_name}]
        )

        existing = self._load_metadata(user_id)
        if any(d["name"] == safe_name for d in existing):
            self.delete_document(safe_name, user_id)

        vectorstore = self._get_vectorstore(user_id)
        vectorstore.add_documents(chunks)

        metadata = self._load_metadata(user_id)
        metadata.append({
            "name": safe_name,
            "chunks": len(chunks),
            "words": word_count,
            "type": "URL",
            "keywords": keywords,
            "uploaded_at": datetime.now(timezone.utc).isoformat()
        })
        self._save_metadata(user_id, metadata)

        return {"chunks": len(chunks), "words": word_count, "filename": safe_name}

    def ingest_text(self, text: str, title: str, user_id: str) -> Dict:
        safe_title = re.sub(r'[^\w.\- ]', '_', title.strip())[:80]
        filename = safe_title + ".txt"

        word_count = len(text.split())
        keywords = self._extract_keywords(text)

        chunks = self.splitter.create_documents(
            texts=[text],
            metadatas=[{"source": filename, "filename": filename}]
        )

        existing = self._load_metadata(user_id)
        if any(d["name"] == filename for d in existing):
            self.delete_document(filename, user_id)

        vectorstore = self._get_vectorstore(user_id)
        vectorstore.add_documents(chunks)

        metadata = self._load_metadata(user_id)
        metadata.append({
            "name": filename,
            "chunks": len(chunks),
            "words": word_count,
            "type": "TXT",
            "keywords": keywords,
            "uploaded_at": datetime.now(timezone.utc).isoformat()
        })
        self._save_metadata(user_id, metadata)

        return {"chunks": len(chunks), "words": word_count, "filename": filename}

    # ── Query helpers ─────────────────────────────────────────────────────────

    def _build_prompt_and_sources(self, question: str, history: List[dict], user_id: str, filter_filename: Optional[str] = None):
        """Shared logic for query() and query_stream(). Returns (prompt, sources, avg_confidence, use_context)."""
        history = history[-10:]

        vectorstore = self._get_vectorstore(user_id)

        try:
            count = len(vectorstore.get()['ids'])
        except Exception:
            count = 0

        if count == 0:
            return None, [], 0.0, False

        search_kwargs = {"k": 4}
        if filter_filename:
            search_kwargs["filter"] = {"filename": filter_filename}

        docs_with_scores = vectorstore.similarity_search_with_score(question, **search_kwargs)

        scored = []
        for doc, score in docs_with_scores:
            confidence = max(0, 1 - (score / 2))
            scored.append((doc, score, confidence))

        avg_confidence = sum(c for _, _, c in scored) / len(scored) if scored else 0
        max_confidence = max((c for _, _, c in scored), default=0)

        RELEVANCE_THRESHOLD = 0.20
        CASUAL_THRESHOLD = 0.10
        use_context = max_confidence >= RELEVANCE_THRESHOLD

        if use_context:
            context = "\n\n".join([doc.page_content for doc, _, _ in scored])
            prompt = f"""You are RAGs_AI, a document assistant. The user has uploaded documents and is asking a question.
Answer using the document context below. Be precise and always state which document the answer comes from.
If the context does not contain enough information, say so clearly — do NOT make up an answer.

Document Context:
{context}

Question: {question}

Answer:"""
        elif max_confidence >= CASUAL_THRESHOLD:
            prompt = f"""You are RAGs_AI, a document assistant. The user has uploaded documents but their question did not closely match anything in them.
Respond helpfully, let them know you couldn't find a strong match in their documents, and suggest they rephrase or be more specific.

Question: {question}

Answer:"""
        else:
            prompt = f"""You are RAGs_AI, a friendly assistant. The user is making small talk. Respond naturally and warmly.

Question: {question}

Answer:"""

        sources = []
        if use_context:
            seen_files = set()
            for doc, _, confidence in scored:
                fname = doc.metadata.get("filename", "unknown")
                if fname not in seen_files and confidence >= CASUAL_THRESHOLD:
                    seen_files.add(fname)
                    sources.append({
                        "file": fname,
                        "page": doc.metadata.get("page", 1),
                        "snippet": doc.page_content[:120].strip() + "..."
                    })
        sources = sources[:2]

        return prompt, sources, round(avg_confidence, 2) if use_context else None, use_context

    # ── Sync query ────────────────────────────────────────────────────────────

    def query(self, question: str, history: List[dict], user_id: str, filter_filename: Optional[str] = None) -> Dict:
        history = history[-10:]
        vectorstore = self._get_vectorstore(user_id)

        try:
            count = len(vectorstore.get()['ids'])
        except Exception:
            count = 0

        if count == 0:
            return {
                "answer": "No documents uploaded yet! Please upload a document first.",
                "sources": [],
                "confidence": 0.0,
                "used_context": False
            }

        search_kwargs = {"k": 4}
        if filter_filename:
            search_kwargs["filter"] = {"filename": filter_filename}

        docs_with_scores = vectorstore.similarity_search_with_score(question, **search_kwargs)

        scored = []
        for doc, score in docs_with_scores:
            confidence = max(0, 1 - (score / 2))
            scored.append((doc, score, confidence))

        avg_confidence = sum(c for _, _, c in scored) / len(scored) if scored else 0
        max_confidence = max((c for _, _, c in scored), default=0)

        RELEVANCE_THRESHOLD = 0.20
        CASUAL_THRESHOLD = 0.10
        use_context = max_confidence >= RELEVANCE_THRESHOLD

        chat_history = []
        for msg in history:
            if msg["role"] == "user":
                chat_history.append(HumanMessage(content=msg["content"]))
            elif msg["role"] == "assistant":
                chat_history.append(AIMessage(content=msg["content"]))

        if use_context:
            context = "\n\n".join([doc.page_content for doc, _, _ in scored])
            prompt = f"""You are RAGs_AI, a document assistant. The user has uploaded documents and is asking a question.
Answer using the document context below. Be precise and always state which document the answer comes from.
If the context does not contain enough information, say so clearly — do NOT make up an answer.

Document Context:
{context}

Question: {question}

Answer:"""
        elif max_confidence >= CASUAL_THRESHOLD:
            prompt = f"""You are RAGs_AI, a document assistant. The user has uploaded documents but their question did not closely match anything in them.
Respond helpfully, let them know you couldn't find a strong match in their documents, and suggest they rephrase or be more specific.

Question: {question}

Answer:"""
        else:
            prompt = f"""You are RAGs_AI, a friendly assistant. The user is making small talk. Respond naturally and warmly.

Question: {question}

Answer:"""

        response = self.llm.invoke(prompt)

        sources = []
        if use_context:
            seen_files = set()
            for doc, _, confidence in scored:
                filename = doc.metadata.get("filename", "unknown")
                if filename not in seen_files and confidence >= CASUAL_THRESHOLD:
                    seen_files.add(filename)
                    sources.append({
                        "file": filename,
                        "page": doc.metadata.get("page", 1),
                        "snippet": doc.page_content[:120].strip() + "..."
                    })

        return {
            "answer": response.content,
            "sources": sources[:2],
            "confidence": round(avg_confidence, 2) if use_context else None,
            "used_context": use_context
        }

    # ── Async streaming query ─────────────────────────────────────────────────

    async def query_stream(
        self,
        question: str,
        history: List[dict],
        user_id: str,
        session_id: Optional[str] = None,
        filter_filename: Optional[str] = None
    ) -> AsyncGenerator[str, None]:
        history = history[-10:]

        vectorstore = self._get_vectorstore(user_id)

        try:
            count = len(vectorstore.get()['ids'])
        except Exception:
            count = 0

        if count == 0:
            no_doc_msg = "No documents uploaded yet! Please upload a document first."
            yield f'data: {json.dumps({"type": "token", "content": no_doc_msg})}\n\n'
            done_payload = {"type": "done", "sources": [], "confidence": 0.0, "used_context": False, "session_id": session_id}
            yield f'data: {json.dumps(done_payload)}\n\n'
            return

        search_kwargs = {"k": 4}
        if filter_filename:
            search_kwargs["filter"] = {"filename": filter_filename}

        docs_with_scores = vectorstore.similarity_search_with_score(question, **search_kwargs)

        scored = []
        for doc, score in docs_with_scores:
            confidence = max(0, 1 - (score / 2))
            scored.append((doc, score, confidence))

        avg_confidence = sum(c for _, _, c in scored) / len(scored) if scored else 0
        max_confidence = max((c for _, _, c in scored), default=0)

        RELEVANCE_THRESHOLD = 0.20
        CASUAL_THRESHOLD = 0.10
        use_context = max_confidence >= RELEVANCE_THRESHOLD

        if use_context:
            context = "\n\n".join([doc.page_content for doc, _, _ in scored])
            prompt = f"""You are RAGs_AI, a document assistant. The user has uploaded documents and is asking a question.
Answer using the document context below. Be precise and always state which document the answer comes from.
If the context does not contain enough information, say so clearly — do NOT make up an answer.

Document Context:
{context}

Question: {question}

Answer:"""
        elif max_confidence >= CASUAL_THRESHOLD:
            prompt = f"""You are RAGs_AI, a document assistant. The user has uploaded documents but their question did not closely match anything in them.
Respond helpfully, let them know you couldn't find a strong match in their documents, and suggest they rephrase or be more specific.

Question: {question}

Answer:"""
        else:
            prompt = f"""You are RAGs_AI, a friendly assistant. The user is making small talk. Respond naturally and warmly.

Question: {question}

Answer:"""

        sources = []
        if use_context:
            seen_files = set()
            for doc, _, confidence in scored:
                fname = doc.metadata.get("filename", "unknown")
                if fname not in seen_files and confidence >= CASUAL_THRESHOLD:
                    seen_files.add(fname)
                    sources.append({
                        "file": fname,
                        "page": doc.metadata.get("page", 1),
                        "snippet": doc.page_content[:120].strip() + "..."
                    })
        sources = sources[:2]

        full_answer = ""
        async for chunk in self.llm.astream(prompt):
            token = chunk.content
            if token:
                full_answer += token
                yield f'data: {json.dumps({"type": "token", "content": token})}\n\n'

        confidence_val = round(avg_confidence, 2) if use_context else None

        # Persist to session if session_id provided
        if session_id:
            try:
                self.save_session_exchange(
                    session_id, user_id, question, full_answer,
                    sources, confidence_val, use_context
                )
            except Exception:
                pass

        done_payload = {
            "type": "done",
            "sources": sources,
            "confidence": confidence_val,
            "used_context": use_context,
            "session_id": session_id
        }
        yield f'data: {json.dumps(done_payload)}\n\n'

    # ── Analytics ─────────────────────────────────────────────────────────────

    def get_analytics(self, user_id: str) -> Dict:
        metadata = self._load_metadata(user_id)
        if not metadata:
            return {
                "total_documents": 0,
                "total_chunks": 0,
                "total_words": 0,
                "documents": [],
                "top_keywords": []
            }

        total_chunks = sum(d["chunks"] for d in metadata)
        total_words = sum(d["words"] for d in metadata)

        combined = Counter()
        for doc in metadata:
            for kw in doc.get("keywords", []):
                combined[kw["word"]] += kw["count"]
        top_keywords = [{"word": w, "count": c} for w, c in combined.most_common(10)]

        return {
            "total_documents": len(metadata),
            "total_chunks": total_chunks,
            "total_words": total_words,
            "documents": [
                {
                    "name": d["name"],
                    "chunks": d["chunks"],
                    "words": d["words"],
                    "type": d["type"],
                    "uploaded_at": d.get("uploaded_at")
                }
                for d in metadata
            ],
            "top_keywords": top_keywords
        }

    # ── Delete / Reset ────────────────────────────────────────────────────────

    def delete_document(self, filename: str, user_id: str):
        vectorstore = self._get_vectorstore(user_id)
        try:
            results = vectorstore.get(where={"filename": filename})
            ids_to_delete = results.get("ids", [])
            if ids_to_delete:
                vectorstore.delete(ids=ids_to_delete)
        except Exception:
            pass

        metadata = self._load_metadata(user_id)
        metadata = [d for d in metadata if d["name"] != filename]
        self._save_metadata(user_id, metadata)

    def reset(self, user_id: str):
        import shutil
        vectorstore = self._get_vectorstore(user_id)
        try:
            vectorstore.delete_collection()
        except Exception:
            pass
        user_dir = os.path.join(self.user_data_dir, user_id)
        if os.path.exists(user_dir):
            shutil.rmtree(user_dir)
